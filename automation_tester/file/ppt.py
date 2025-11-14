import asyncio
import io
import logging
from pathlib import Path

from PIL import Image

from automation_tester.file.base_file import BaseFile
from automation_tester.utils.image_parser import parse_image_with_llm

logger = logging.getLogger(__name__)


class PPTFile(BaseFile):
    """
    PPT 文件解析类，支持文本和图片提取

    Args:
        source: 文件路径
    """

    def __init__(self, source: str):
        super().__init__(source)

    async def parse_text(self, **kwargs):
        if not self._path:
            raise ValueError("path is not set")

        file_path = Path(self._path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {self._path}")

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as e:
            logger.error(f"缺少python-pptx依赖: {e}")
            raise ImportError("请安装python-pptx") from e

        try:
            prs = Presentation(self._path)
        except Exception as e:
            logger.error(f"无法打开 PPT 文件: {e}")
            raise RuntimeError(f" PPT 文件格式错误或损坏: {e}") from e

        for slide_index, slide in enumerate(prs.slides, 1):
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        for line in shape.text.strip().splitlines():
                            if line.strip():
                                yield line.strip()

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        if image and image.blob:
                            with Image.open(io.BytesIO(image.blob)) as pil_image:
                                pil_image.thumbnail((1024, 1024))

                                max_retries = 3
                                for attempt in range(max_retries):
                                    try:
                                        result = await parse_image_with_llm(pil_image)
                                        yield result
                                        break
                                    except Exception as e:
                                        if attempt < max_retries - 1:
                                            await asyncio.sleep(1)
                                        else:
                                            logger.warning(f"图片解析失败: {e}")
                                        break

            except Exception as e:
                logger.warning(f"处理 PPT {slide_index} 时出错: {e}")
                continue
